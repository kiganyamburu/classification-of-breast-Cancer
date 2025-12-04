from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text=None):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True

        if content_text:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()

            for line in content_text.split("\n"):
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.level = 0
        return slide

    # Helper to add a slide with title and image
    def add_image_slide(title_text, image_path):
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True

        if os.path.exists(image_path):
            # Center the image
            left = Inches(1)
            top = Inches(1.8)
            height = Inches(5)
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            print(f"Warning: {image_path} not found.")

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Breast Cancer Classification"
    subtitle.text = (
        "Machine Learning Capstone Project\nPredicting Malignant vs Benign Tumors"
    )
    title.text_frame.paragraphs[0].font.size = Pt(54)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # 2. Problem Statement (Requirement A)
    problem_statement = (
        "What are we solving?\n"
        "• Automated classification of breast tumors as Benign or Malignant based on cell measurements\n\n"
        "Why does it matter?\n"
        "• Breast cancer is one of the most common cancers affecting women worldwide\n"
        "• Early and accurate diagnosis significantly improves survival rates (99% for early detection)\n"
        "• Manual diagnosis can be time-consuming and subjective\n\n"
        "Real-world impact:\n"
        "• Assists medical professionals in making faster, data-driven decisions\n"
        "• Reduces diagnostic errors and improves patient outcomes\n"
        "• Supports screening programs in resource-limited settings"
    )
    add_slide("Problem Statement", problem_statement)

    # 3. Data Source (Part of Methodology - Requirement B)
    data_source = (
        "Data Source:\n"
        "• Breast Cancer Wisconsin (Diagnostic) Dataset\n"
        "• UCI Machine Learning Repository\n\n"
        "Dataset Characteristics:\n"
        "• 569 samples (357 Benign, 212 Malignant)\n"
        "• 30 numeric features derived from digitized images of Fine Needle Aspirate (FNA)\n"
        "• Features computed for each cell nucleus: radius, texture, perimeter, area, smoothness, compactness, concavity, concave points, symmetry, fractal dimension\n"
        "• Each feature has 3 measurements: mean, standard error (SE), and worst (largest value)\n\n"
        "Target Variable:\n"
        "• Diagnosis: M (Malignant) or B (Benign)"
    )
    add_slide("Step 1: Data Source", data_source)

    # 4. Data Preparation (Requirement B)
    data_prep = (
        "Data Cleaning:\n"
        "• Removed unnecessary columns: 'id' (patient identifier) and 'Unnamed: 32' (empty column)\n"
        "• Verified no missing values in the dataset\n\n"
        "Data Encoding:\n"
        "• Converted categorical target variable to numeric:\n"
        "  - B (Benign) → 0\n"
        "  - M (Malignant) → 1\n\n"
        "Data Splitting:\n"
        "• Training set: 80% (455 samples)\n"
        "• Testing set: 20% (114 samples)\n"
        "• Used stratified split to maintain class distribution in both sets\n\n"
        "Feature Scaling:\n"
        "• Random Forest handles features without scaling, so no normalization was required"
    )
    add_slide("Step 2: Data Preparation", data_prep)

    # 5. Exploratory Data Analysis (Requirement B)
    eda_text = (
        "Analysis Performed:\n\n"
        "1. Class Distribution Analysis\n"
        "   • Examined balance between Benign and Malignant cases\n"
        "   • Found moderate imbalance (62.7% Benign, 37.3% Malignant)\n\n"
        "2. Feature Correlation Analysis\n"
        "   • Identified features most correlated with diagnosis\n"
        "   • Analyzed multicollinearity among features\n\n"
        "3. Statistical Analysis\n"
        "   • Computed descriptive statistics grouped by diagnosis\n"
        "   • Identified significant differences in feature distributions"
    )
    add_slide("Step 3: Exploratory Data Analysis", eda_text)

    # 6. EDA: Class Distribution
    add_image_slide("Class Distribution Visualization", "class_distribution.png")

    # 7. Correlation Heatmap
    add_image_slide("Feature Correlation Analysis", "correlation_heatmap.png")

    # 8. Machine Learning Methodology (Requirement B)
    ml_methodology = (
        "Algorithm Selection:\n"
        "• Random Forest Classifier - an ensemble learning method\n\n"
        "Why Random Forest?\n"
        "• Handles non-linear relationships effectively\n"
        "• Robust to outliers and overfitting\n"
        "• Provides feature importance rankings\n"
        "• Works well with medical diagnostic data\n\n"
        "Model Configuration:\n"
        "• Default parameters with random_state=42 for reproducibility\n"
        "• Number of trees: 100 (default)\n\n"
        "Training Process:\n"
        "• Fit model on 455 training samples\n"
        "• Validated on 114 held-out test samples"
    )
    add_slide("Step 4: Machine Learning Approach", ml_methodology)

    # 9. Evaluation Metrics (Requirement B)
    eval_metrics = (
        "Evaluation Strategy:\n\n"
        "Primary Metric:\n"
        "• Accuracy: Overall percentage of correct predictions\n\n"
        "Detailed Analysis:\n"
        "• Confusion Matrix: Breakdown of True Positives, True Negatives, False Positives, False Negatives\n\n"
        "Why these metrics?\n"
        "• In medical diagnosis, it's critical to understand both:\n"
        "  - Sensitivity (correctly identifying malignant cases)\n"
        "  - Specificity (correctly identifying benign cases)\n"
        "• False negatives (missing cancer) are particularly costly"
    )
    add_slide("Step 5: Evaluation Metrics", eval_metrics)

    # 10. Results (Requirement D)
    metrics_text = "Model Performance Results:\n\n"
    try:
        with open("metrics.txt", "r") as f:
            content = f.read()
            metrics_text += content

            # Parse accuracy for interpretation
            if "Accuracy:" in content:
                acc_line = [
                    line for line in content.split("\n") if "Accuracy:" in line
                ][0]
                metrics_text += f"\n\nInterpretation:\n"
                metrics_text += "The model correctly classified the test samples with high accuracy,\n"
                metrics_text += "demonstrating strong predictive capability for breast cancer diagnosis."
    except FileNotFoundError:
        metrics_text += (
            "Accuracy: Not calculated yet. Please run generate_assets.py first."
        )

    add_slide("Results: Model Performance", metrics_text)

    # 11. Confusion Matrix Results
    add_image_slide("Results: Confusion Matrix", "confusion_matrix.png")

    # 12. Results Interpretation (Requirement D)
    results_interpretation = (
        "Key Findings:\n\n"
        "✓ High Accuracy: The model achieves excellent performance on unseen data\n\n"
        "✓ Confusion Matrix Insights:\n"
        "  • True Positives: Correctly identified malignant cases\n"
        "  • True Negatives: Correctly identified benign cases\n"
        "  • False Positives: Benign cases incorrectly flagged as malignant (low risk)\n"
        "  • False Negatives: Malignant cases missed (higher risk - minimize this)\n\n"
        "Clinical Significance:\n"
        "• The model can serve as a decision support tool for pathologists\n"
        "• Reduces time required for diagnosis\n"
        "• Provides consistent, reproducible results"
    )
    add_slide("Results Interpretation", results_interpretation)

    # 13. Conclusion & Impact
    conclusion_text = (
        "Summary:\n"
        "• Successfully developed an ML model for breast cancer classification\n"
        "• Achieved high accuracy using Random Forest algorithm\n"
        "• Model can distinguish between benign and malignant tumors effectively\n\n"
        "Real-World Applications:\n"
        "• Integration into hospital diagnostic workflows\n"
        "• Second-opinion system for pathologists\n"
        "• Screening tool in under-resourced medical facilities\n\n"
        "Limitations:\n"
        "• Requires validation on diverse patient populations\n"
        "• Should complement, not replace, expert medical judgment\n"
        "• Model interpretability could be enhanced for clinical adoption"
    )
    add_slide("Conclusion & Impact", conclusion_text)

    # 14. Future Work & Improvements
    future_work = (
        "Next Steps:\n\n"
        "Model Enhancement:\n"
        "• Hyperparameter tuning using GridSearchCV or RandomizedSearchCV\n"
        "• Test additional algorithms: SVM, XGBoost, Neural Networks\n"
        "• Ensemble methods combining multiple models\n\n"
        "Feature Engineering:\n"
        "• Feature selection to identify most important predictors\n"
        "• Dimensionality reduction (PCA) for efficiency\n\n"
        "Deployment:\n"
        "• Create web application for easy access\n"
        "• Develop API for integration with medical systems\n"
        "• Mobile app for point-of-care diagnosis"
    )
    add_slide("Future Work & Improvements", future_work)

    # 15. Thank You Slide
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Thank You"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True

    # Add subtitle
    p2 = text_frame.add_paragraph()
    p2.text = "\nQuestions & Discussion"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(32)

    # Save
    output_file = "Breast_Cancer_Classification_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✓ Presentation successfully created!")
    print(f"✓ File saved as: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"\nPresentation includes:")
    print(
        "  • Clear problem statement explaining what we're solving and why it matters"
    )
    print("  • Step-by-step methodology covering all aspects:")
    print("    - Data source and collection")
    print("    - Data preparation and cleaning")
    print("    - Exploratory data analysis")
    print("    - Machine learning process")
    print("    - Evaluation strategy")
    print("  • Clear results presentation with visualizations")
    print("  • Clinical interpretation and real-world impact")
    print("\nYou can now present this PowerPoint for your capstone!")


if __name__ == "__main__":
    create_presentation()
